#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI 管理學 - 培訓簡報生成器
使用 python-pptx 函式庫生成 52 頁簡報
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 建立簡報
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 顏色定義
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)  # 封面背景
TITLE_BLUE = RGBColor(0x16, 0x21, 0x3e)  # 標題深藍
ACCENT_BLUE = RGBColor(0x0f, 0x34, 0x6a)  # 強調藍
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)

current_page = 0


def add_page_number(slide, page_num):
    """添加頁碼"""
    left = Inches(12.0)
    top = Inches(7.0)
    width = Inches(1.0)
    height = Inches(0.3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY


def add_cover_slide(title, subtitle, author):
    """創建封面頁"""
    global current_page
    slide_layout = prs.slide_layouts[6]  # 空白版面
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # 裝飾線條
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.8), Inches(10.333), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    line.line.fill.background()
    
    # 主標題
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # 副標題
    top = Inches(4.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    
    # 作者
    top = Inches(6.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = author
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    
    current_page += 1


def add_title_slide(title):
    """創建章節標題頁"""
    global current_page
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 淺色背景區塊
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xf5, 0xf7, 0xfa)
    bg.line.fill.background()
    
    # 標題
    left = Inches(0.8)
    top = Inches(0.8)
    width = Inches(11.733)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TITLE_BLUE
    
    # 底線裝飾
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(2), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    
    current_page += 1
    add_page_number(slide, current_page)


def add_content_slide(title, bullets):
    """創建內容頁"""
    global current_page
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 頂部裝飾條
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ACCENT_BLUE
    header_bar.line.fill.background()
    
    # 標題
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(11.733)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_BLUE
    
    # 內容項目
    top = Inches(1.6)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.space_after = Pt(18)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = BLACK
    
    current_page += 1
    add_page_number(slide, current_page)


def add_toc_slide():
    """創建目錄頁"""
    global current_page
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 頂部裝飾條
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ACCENT_BLUE
    header_bar.line.fill.background()
    
    # 標題
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(11.733)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "目錄"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TITLE_BLUE
    
    # 章節列表
    chapters = [
        "第 1 章：把 AI 當人——管理思維的轉換",
        "第 2 章：六大核心管理框架",
        "第 3 章：五種任務類型與 BCG 矩陣",
        "第 4 章：O + KR 指令設計",
        "第 5 章：AI 等級分類與 AI 管理 AI",
        "第 6 章：上班族的工具組合",
        "第 7 章：三層次工作流",
        "第 8 章：分層管理架構",
        "第 9 章：四維度審核",
        "第 10 章：三條線與變革管理",
        "第 11 章：雙軌 PDCA",
        "第 12 章：90 天計劃與核心心法"
    ]
    
    # 左欄
    top = Inches(1.6)
    txBox = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(5.2))
    tf = txBox.text_frame
    for i, ch in enumerate(chapters[:6]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ch
        p.space_after = Pt(16)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = BLACK
    
    # 右欄
    txBox2 = slide.shapes.add_textbox(Inches(6.8), top, Inches(5.5), Inches(5.2))
    tf2 = txBox2.text_frame
    for i, ch in enumerate(chapters[6:]):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = ch
        p.space_after = Pt(16)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = BLACK
    
    current_page += 1
    add_page_number(slide, current_page)


def add_ending_slide():
    """創建結束頁"""
    global current_page
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # 謝謝文字
    left = Inches(0.5)
    top = Inches(3.0)
    width = Inches(12.333)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "謝謝"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # 副標題
    top = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "開始你的 AI 管理之旅"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    
    current_page += 1


# ==================== 開始生成簡報 ====================

# 封面頁
add_cover_slide("AI 管理學", "用管理框架馭 AI 能力", "草台班子研究室")

# 目錄頁
add_toc_slide()

# ==================== 第 1 章 ====================
add_title_slide("第 1 章：把 AI 當人——管理思維的轉換")

add_content_slide("把 AI 當人——管理思維的轉換", [
    "傳統工具思維：輸入 → 輸出，一成不變",
    "AI 管理思維：目標 → 溝通 → 迭代優化",
    "把 AI 視為團隊成員，而非冷冰冰的工具",
    "需要明確的目標、清楚的指令、持續的反饋",
    "關鍵轉變：從「使用」到「管理」"
])

add_content_slide("AI 像一個下屬或助理", [
    "優點：快速、不知疲倦、知識淵博",
    "可處理大量重複性、知識性工作",
    "協助腦力激盪、資料整理、初稿撰寫",
    "如同一位聰明但需要指導的助理",
    "需要你的方向感與判斷力"
])

add_content_slide("AI 三大缺陷", [
    "沒有 Common Sense（常識）",
    "    可能給出邏輯上合理但實際荒謬的答案",
    "容易幻覺（Hallucination）",
    "    編造不存在的資訊、數據或引用",
    "不知道你的具體情況",
    "    缺乏對你公司、產業、個人背景的深度理解",
    "應對策略：審核、驗證、提供上下文"
])

add_content_slide("真實失敗案例", [
    "律師事務所引用 AI 生成的虛構判例",
    "公司使用 AI 產生的錯誤財務數據",
    "學生提交含有虛構參考文獻的論文",
    "共同點：盲目信任，缺乏驗證",
    "教訓：AI 輸出必須經過人工審核"
])

# ==================== 第 2 章 ====================
add_title_slide("第 2 章：六大核心管理框架")

add_content_slide("六大核心框架", [
    "OKR：目標與關鍵結果（Objectives & Key Results）",
    "MECE：相互獨立，完全窮盡（Mutually Exclusive, Collectively Exhaustive）",
    "BCG 矩陣：波士頓顧問群矩陣",
    "SWOT：優勢、劣勢、機會、威脅分析",
    "80-20 法則：帕累托法則，抓重點",
    "STAR：情境、任務、行動、結果"
])

add_content_slide("OKR × AI：目標導向的指令設計", [
    "Objective（目標）：你要達成什麼？",
    "Key Results（關鍵結果）：如何衡量成功？",
    "將 OKR 應用於 Prompt 設計",
    "範例：寫一份報告 → 目標：說服高層投資",
    "明確的 KR：涵蓋市場規模、ROI、風險評估"
])

add_content_slide("MECE × AI：結構化思考", [
    "確保思考無遺漏、不重疊",
    "用於分解複雜問題",
    "常見分類方式：",
    "    • 內部 / 外部",
    "    • 短期 / 長期",
    "    • 產品 / 市場 / 營運",
    "讓 AI 輸出結構更清晰完整"
])

add_content_slide("BCG 矩陣 × AI：任務分類與資源配置", [
    "明星（Stars）：高成長、高市佔 → 重點投入",
    "金牛（Cash Cows）：低成長、高市佔 → 維持運作",
    "問號（Question Marks）：高成長、低市佔 → 評估潛力",
    "瘦狗（Dogs）：低成長、低市佔 → 考慮淘汰",
    "應用於 AI 任務優先級排序"
])

# ==================== 第 3 章 ====================
add_title_slide("第 3 章：五種任務類型與 BCG 矩陣")

add_content_slide("五種任務類型", [
    "創意生成型：腦力激盪、文案發想、概念設計",
    "分析整理型：資料彙整、數據分析、摘要重點",
    "撰寫潤稿型：郵件、報告、文章、翻譯",
    "程式開發型：代碼編寫、除錯、技術文件",
    "決策輔助型：方案比較、風險評估、建議提供"
])

add_content_slide("BCG 任務分配矩陣", [
    "高價值 + 高可行性 = 立即執行（AI 主力）",
    "高價值 + 低可行性 = 拆解或人工介入",
    "低價值 + 高可行性 = 自動化批次處理",
    "低價值 + 低可行性 = 暫緩或放棄",
    "根據任務特性選擇適合的 AI 工具與策略"
])

# ==================== 第 4 章 ====================
add_title_slide("第 4 章：O + KR 指令結構")

add_content_slide("O + KR 指令結構", [
    "O（Objective）：清楚說明目標",
    "    「我需要一份市場分析報告」",
    "KR（Key Results）：定義成功標準",
    "    • 包含 3 個競爭對手分析",
    "    • 提供具體數據支持",
    "    • 結論需有可行建議",
    "組合使用，讓 AI 精準理解需求"
])

add_content_slide("四種指令模板", [
    "商業報告模板：背景 + 目標 + 格式 + 字數限制",
    "創意發想模板：主題 + 風格 + 數量 + 限制條件",
    "郵件撰寫模板：對象 + 目的 + 語氣 + 行動呼籲",
    "翻譯潤稿模板：原文 + 目標語言 + 風格 + 專業領域"
])

# ==================== 第 5 章 ====================
add_title_slide("第 5 章：AI 等級分類與 AI 管理 AI")

add_content_slide("AI 等級分類", [
    "旗艦級（Flagship）：GPT-4、Claude 3 Opus 等",
    "    複雜推理、創意、長文本處理",
    "主力級（Main）：GPT-4o、Claude 3.5 Sonnet 等",
    "    日常任務、快速回應、性價比高",
    "專門級（Specialized）：代碼、設計、法律等專業模型",
    "輔助級（Assist）：輕量任務、快速查詢、簡單生成"
])

add_content_slide("用 AI 管理 AI", [
    "Prompt 生成器：用 AI 幫你寫更好的 Prompt",
    "Multi-Agent 架構：多個 AI 協作完成任務",
    "    • 研究員 → 收集資訊",
    "    • 分析師 → 整理分析",
    "    • 撰寫者 → 產出報告",
    "AI 審核：用另一個 AI 檢查輸出品質"
])

# ==================== 第 6 章 ====================
add_title_slide("第 6 章：上班族的工具組合")

add_content_slide("上班族的工具組合", [
    "通用對話：ChatGPT、Claude、Gemini",
    "寫作輔助：Notion AI、Grammarly、Copilot",
    "簡報製作：Gamma、Beautiful.ai、Tome",
    "數據分析：Julius AI、ChatCSV、Ajelix",
    "會議記錄：Otter.ai、Fireflies、Fathom",
    "專案管理：Notion、Asana AI、Monday.com"
])

add_content_slide("五個選擇維度", [
    "任務類型：創意 / 分析 / 撰寫 / 程式？",
    "輸出品質：需要最高品質還是夠用即可？",
    "成本考量：免費、訂閱、按量計費？",
    "整合需求：是否需要與現有工具整合？",
    "安全隱私：資料是否可上傳雲端？"
])

# ==================== 第 7 章 ====================
add_title_slide("第 7 章：三層次工作流")

add_content_slide("三層次工作流", [
    "第一層：模板（Template）",
    "    標準化的 Prompt 框架，可重複使用",
    "第二層：工作流（Workflow）",
    "    串連多個 AI 任務的自動化流程",
    "第三層：反饋循環（Feedback Loop）",
    "    持續收集結果、優化模板、迭代改進"
])

add_content_slide("CLEAR 框架", [
    "C - Context（上下文）：提供背景資訊",
    "L - Limit（限制）：設定邊界條件",
    "E - Example（範例）：給出具體例子",
    "A - Action（行動）：明確要求什麼",
    "R - Review（審查）：規劃檢查機制",
    "五個步驟，打造高品質 AI 協作流程"
])

# ==================== 第 8 章 ====================
add_title_slide("第 8 章：分層管理架構")

add_content_slide("分層管理架構", [
    "系統層（System Level）：",
    "    AI 使用政策、資安規範、成本管控",
    "專案層（Project Level）：",
    "    專案專用的 AI 工作流、模板、SOP",
    "任務層（Task Level）：",
    "    單一任務的 Prompt、執行、審核"
])

add_content_slide("AI 入職培訓三階段", [
    "第一階段：認識 AI（1-2 週）",
    "    了解能力邊界、嘗試各種應用",
    "第二階段：熟練工具（3-4 週）",
    "    掌握 Prompt 技巧、建立個人模板",
    "第三階段：整合工作（5-8 週）",
    "    將 AI 融入日常工作流程、持續優化"
])

# ==================== 第 9 章 ====================
add_title_slide("第 9 章：四維度審核")

add_content_slide("四維度審核", [
    "準確性（Accuracy）：事實是否正確？",
    "相關性（Relevance）：是否符合需求？",
    "完整性（Completeness）：有無遺漏重要資訊？",
    "可執行性（Actionability）：能否實際執行？",
    "每個維度 1-5 分，建立評分標準"
])

add_content_slide("AI 檢查 AI", [
    "使用不同的 AI 進行交叉驗證",
    "專門的審核 Prompt 設計",
    "請 AI 列出潛在錯誤或遺漏",
    "對比多個 AI 的輸出差異",
    "重要決策仍需人工最終確認"
])

# ==================== 第 10 章 ====================
add_title_slide("第 10 章：三條線與變革管理")

add_content_slide("三條線", [
    "資訊分級：",
    "    公開 / 內部 / 機密 / 絕密",
    "外流邊界：",
    "    哪些資料可以輸入 AI？哪些不行？",
    "管理 AI：",
    "    使用紀錄、審查機制、責任歸屬"
])

add_content_slide("變革管理：員工對 AI 的三種態度", [
    "擁抱者（Enthusiasts）：積極學習，主動嘗試",
    "    → 給予資源，讓他們成為種子",
    "觀望者（Observers）：等待證明，謹慎評估",
    "    → 提供成功案例，降低進入門檻",
    "抗拒者（Resistors）：擔心取代，抗拒改變",
    "    → 強調協作而非取代，提供培訓支持"
])

# ==================== 第 11 章 ====================
add_title_slide("第 11 章：雙軌 PDCA")

add_content_slide("雙軌 PDCA", [
    "Prompt PDCA：",
    "    Plan（計劃）：設計 Prompt",
    "    Do（執行）：輸入 AI 取得結果",
    "    Check（檢查）：評估輸出品質",
    "    Act（行動）：優化 Prompt",
    "AI 策略 PDCA：",
    "    整體 AI 導入策略的持續改進"
])

# ==================== 第 12 章 ====================
add_title_slide("第 12 章：90 天計劃與核心心法")

add_content_slide("30 天計劃", [
    "第 1-10 天：建立基礎",
    "    • 選定 1-2 個主力工具",
    "    • 建立個人 Prompt 模板庫",
    "第 11-20 天：擴大應用",
    "    • 將 AI 融入 3-5 個日常工作場景",
    "    • 建立部門級工作流",
    "第 21-30 天：優化迭代",
    "    • 評估成效、調整策略",
    "    • 建立長期優化機制"
])

add_content_slide("核心心法：標準比工具重要", [
    "沒有標準，再好的工具也會用錯",
    "建立 Prompt 審核標準",
    "建立輸出品質標準",
    "建立 AI 使用行為標準",
    "持續迭代標準，而非追逐新工具"
])

# ==================== 討論時間 ====================
add_title_slide("討論時間")

add_content_slide("互動問題", [
    "在你的工作中，哪個任務最適合交給 AI 處理？",
    "    思考：價值高 + 可行性高的任務",
    "你會如何設計一個 AI 入職培訓計劃？",
    "    思考：三階段培訓的具體內容",
    "如何建立團隊的 AI 輸出品質標準？",
    "    思考：四維度審核的落地執行"
])

# ==================== 行動清單 ====================
add_title_slide("行動清單")

add_content_slide("5 個具體行動項目", [
    "本週：選定 1 個主力 AI 工具，完成註冊與基礎設定",
    "本週：建立 3 個個人專用 Prompt 模板",
    "本月：找出 3 個可交給 AI 處理的重複性任務",
    "本月：與團隊分享一個 AI 使用成功案例",
    "本季：建立部門級的 AI 使用 SOP 與審核標準"
])

# 結束頁
add_ending_slide()

# 儲存簡報
output_path = "/Volumes/WorkData/ctbzai/docs/ai-management-book/training-deck/AI管理學培訓簡報.pptx"
prs.save(output_path)

print(f"✅ 簡報已生成：{output_path}")
print(f"📊 總頁數：{current_page} 頁")
print("\n簡報結構：")
print("  - 封面頁")
print("  - 目錄頁")
print("  - 第 1-12 章內容（各章節標題頁 + 內容頁）")
print("  - 討論時間")
print("  - 行動清單")
print("  - 結束頁")
